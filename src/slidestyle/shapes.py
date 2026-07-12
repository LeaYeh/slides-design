"""Shape kit — the locked-style architecture/diagram building blocks."""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from . import tokens as t

def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color

def _line(shape, color, width_pt=1.5):
    shape.line.color.rgb = color; shape.line.width = Pt(width_pt)

def _no_line(shape):
    shape.line.fill.background()

def box(slide, pos, text, *, kind="plain", accent="blue", sub=None):
    """pos=(left,top,width,height) in inches. kind: plain|focus|neutral."""
    left, top, width, height = (Inches(v) for v in pos)
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sp.adjustments[0] = 0.08  # ~10px radius
    if kind == "focus":
        _fill(sp, t.ramp(accent, "c100")); _line(sp, t.ramp(accent, "c600"), 1.75)
        txt_color = t.ramp(accent, "c800")
    elif kind == "neutral":
        _fill(sp, t.MIST); _no_line(sp); txt_color = t.GRAPHITE
    else:  # plain
        _fill(sp, t.WHITE); _line(sp, t.HAIRLINE, 1.5); txt_color = t.GRAPHITE
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = "Helvetica Neue"; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = txt_color
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub
        r2.font.name = "Helvetica Neue"; r2.font.size = Pt(11); r2.font.color.rgb = t.SLATE
    return sp

def arrow(slide, x1, y1, x2, y2):
    """Thin connector with an arrowhead, Silver."""
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = t.SILVER; cn.line.width = Pt(1.5)
    # add end arrowhead via XML
    ln = cn.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return cn

def band(slide, pos, label, tag, accent):
    """Horizontal layer band (Route A pattern B)."""
    left, top, width, height = (Inches(v) for v in pos)
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sp.adjustments[0] = 0.06
    _fill(sp, t.ramp(accent, "c100")); _line(sp, t.ramp(accent, "c600"), 1.75)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "   " + label
    r.font.name = "Helvetica Neue"; r.font.size = Pt(16); r.font.bold = True
    r.font.color.rgb = t.ramp(accent, "c800")
    if tag:
        # small uppercase category tag, right-aligned inside the band
        tb = slide.shapes.add_textbox(left + width - Inches(2.2), top, Inches(2.0), height)
        ttf = tb.text_frame; ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tp = ttf.paragraphs[0]; tp.alignment = PP_ALIGN.RIGHT
        tr = tp.add_run(); tr.text = tag.upper()
        tr.font.name = "Helvetica Neue"; tr.font.size = Pt(11); tr.font.bold = True
        tr.font.color.rgb = t.ramp(accent, "c800")
        tr.font._rPr.set("spc", str(int(t.KICKER_TRACK_PT * 100)))
    return sp

def kpi_tile(slide, pos, number, label, *, accent=None):
    """Color card. accent=None -> neutral Mist tile; else tint tile with accent number."""
    left, top, width, height = (Inches(v) for v in pos)
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sp.adjustments[0] = 0.06
    _fill(sp, t.ramp(accent, "c100") if accent else t.MIST); _no_line(sp)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = number
    r.font.name = "Helvetica Neue"; r.font.size = Pt(30); r.font.bold = True
    r.font.color.rgb = t.ramp(accent, "c800") if accent else t.INK
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = label
    r2.font.name = "Helvetica Neue"; r2.font.size = Pt(12); r2.font.color.rgb = t.SLATE
    return sp

# --- swim-lane / composition kit (for the layout-rich architecture example) ---

def _centered(shape, text, *, size=13, bold=True, color=t.GRAPHITE):
    tf = shape.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = "Helvetica Neue"; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return shape

def _small_caps(slide, pos, text, color, *, align=PP_ALIGN.LEFT, size=11):
    left, top, width, height = (Inches(v) for v in pos)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text.upper()
    r.font.name = "Helvetica Neue"; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = color
    r.font._rPr.set("spc", str(int(t.KICKER_TRACK_PT * 100)))
    return tb

def container(slide, pos, title, *, accent=None):
    """Titled subsystem group box: Mist surface, small-caps title top-left."""
    left, top, width, height = (Inches(v) for v in pos)
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sp.adjustments[0] = 0.03
    _fill(sp, t.MIST); _line(sp, t.HAIRLINE, 1.25)
    _small_caps(slide, (pos[0] + 0.22, pos[1] + 0.12, pos[2] - 0.44, 0.34),
                title, t.ramp(accent, "c800") if accent else t.SLATE)
    return sp

def cylinder(slide, pos, label, *, accent=None):
    """Data-store cylinder (a shared resource, often straddling lanes)."""
    left, top, width, height = (Inches(v) for v in pos)
    sp = slide.shapes.add_shape(MSO_SHAPE.CAN, left, top, width, height)
    _fill(sp, t.ramp(accent, "c100") if accent else t.MIST); _line(sp, t.HAIRLINE, 1.5)
    _centered(sp, label, size=11, color=t.GRAPHITE)
    return sp

def document(slide, pos, label, *, accent="blue"):
    """Artifact / output — flowchart document shape."""
    left, top, width, height = (Inches(v) for v in pos)
    sp = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_DOCUMENT, left, top, width, height)
    _fill(sp, t.ramp(accent, "c100")); _line(sp, t.ramp(accent, "c600"), 1.5)
    _centered(sp, label, size=11, color=t.ramp(accent, "c800"))
    return sp

def diamond(slide, pos, label, *, accent="blue"):
    """Decision node."""
    left, top, width, height = (Inches(v) for v in pos)
    sp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, width, height)
    _fill(sp, t.WHITE); _line(sp, t.ramp(accent, "c600"), 1.5)
    _centered(sp, label, size=10, color=t.ramp(accent, "c800"))
    return sp

def lane_label(slide, center_y, text):
    """Vertical (rotated) environment label on the far left gutter."""
    w = 2.2  # unrotated width spans the lane height after 270deg rotation
    tb = slide.shapes.add_textbox(Inches(0.30 - w / 2), Inches(center_y - 0.25), Inches(w), Inches(0.5))
    tb.rotation = 270
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text.upper()
    r.font.name = "Helvetica Neue"; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = t.SLATE
    r.font._rPr.set("spc", str(int(t.KICKER_TRACK_PT * 100)))
    return tb

_LINE = t.rgb("B4B7BD")  # refined connector gray

def _seg(slide, x1, y1, x2, y2, color, *, arrow=False, width=1.0, head="sm"):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(width)
    if arrow:
        from pptx.oxml.ns import qn
        ln = cn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": head, "len": head}))
    return cn

def _chip(slide, cx, cy, text):
    """A small caption on a white background so lines don't cross it."""
    w = len(text) * 0.075 + 0.18
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx - w / 2), Inches(cy - 0.13),
                               Inches(w), Inches(0.26))
    _fill(r, t.WHITE); _no_line(r)
    _small_caps(slide, (cx - w / 2, cy - 0.13, w, 0.26), text, t.SLATE,
                align=PP_ALIGN.CENTER, size=9)

def elbow(slide, pts, *, label=None, color=None):
    """Orthogonal connector through pts (right-angle bends), arrowhead on the last
    segment. pts = [(x,y), ...] in inches. Optional label chip on the first segment."""
    color = color or _LINE
    for i in range(len(pts) - 1):
        (x1, y1), (x2, y2) = pts[i], pts[i + 1]
        _seg(slide, x1, y1, x2, y2, color, arrow=(i == len(pts) - 2))
    if label:
        (x1, y1), (x2, y2) = pts[0], pts[1]
        _chip(slide, (x1 + x2) / 2, (y1 + y2) / 2, label)

def feedback_arc(slide, x1, x2, y_box_top, label, *, rise=0.4):
    """Loop-back indicator above a row: up from x1, across, down to x2 (arrowhead)."""
    y_arc = y_box_top - rise
    _seg(slide, x1, y_box_top, x1, y_arc, _LINE)
    _seg(slide, x1, y_arc, x2, y_arc, _LINE)
    _seg(slide, x2, y_arc, x2, y_box_top, _LINE, arrow=True)
    lo, hi = sorted((x1, x2))
    _small_caps(slide, (lo, y_arc - 0.32, hi - lo, 0.3), label, t.SLATE,
                align=PP_ALIGN.CENTER, size=10)
