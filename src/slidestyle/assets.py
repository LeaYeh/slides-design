"""Personal assets for architecture diagrams — photo + tool icons.

Source of truth for the photo + core skills is the jsonresume registry
(https://registry.jsonresume.org/LeaYeh). Tool icons are additionally driven by the
private JD tracker in the local lyeh-infra repo (`docs/resume/jd-tracker.md`) so that
logos for target-role tech (Spark, Kafka, Terraform, AWS, …) are ready too.

`sync()` fetches the resume, downloads the profile photo, reads whichever source docs
are available, and downloads an official-color logo for every recognised tool that
appears in them — caching everything under `assets/` (committed). All reads go through
the local manifest, so decks build offline.

The JD tracker is gitignored/private; only generic brand-logo *names* land in the
committed manifest — no company names, salaries, or other JD content is ever stored.
If the lyeh-infra repo isn't present locally, the JD tools are simply skipped.

Icon colors are official brand logos: Devicon (true multi-color) where available,
Simple Icons (single official-color silhouette) for tools Devicon lacks.

    python -m slidestyle.assets            # refresh the local cache from the sources
"""
import base64
import json
import os
import re
import urllib.request
from pathlib import Path

REGISTRY_URL = "https://registry.jsonresume.org/LeaYeh.json"
ROOT = Path(__file__).resolve().parents[2]        # repo root (…/slides)
ASSETS = ROOT / "assets"
MANIFEST = ASSETS / "manifest.json"

# Private, gitignored JD tracker in the sibling lyeh-infra repo. Override with
# SLIDESTYLE_JD_DOC; skipped silently if absent (keeps this repo self-contained).
JD_DOC = Path(os.environ.get(
    "SLIDESTYLE_JD_DOC",
    ROOT.parent / "git_dev" / "lyeh-infra" / "docs" / "resume" / "jd-tracker.md"))

_DEVICON = "https://cdn.jsdelivr.net/gh/devicons/devicon@latest/icons/{slug}.svg"
_SIMPLE = "https://cdn.simpleicons.org/{slug}/{color}"

# name -> (source, slug, color|None, [lowercase aliases matched against CV keywords])
#   devicon slug = "<dir>/<file-without-.svg>"  · simpleicons slug = brand id, color = hex6
ICON_MAP = {
    "Docker":          ("devicon", "docker/docker-original", None, ["docker"]),
    "Kubernetes":      ("devicon", "kubernetes/kubernetes-plain", None, ["kubernetes", "k3s", "k8s"]),
    "PyTorch":         ("devicon", "pytorch/pytorch-original", None, ["pytorch"]),
    "Prometheus":      ("devicon", "prometheus/prometheus-original", None, ["prometheus"]),
    "Grafana":         ("devicon", "grafana/grafana-original", None, ["grafana"]),
    "Apache Airflow":  ("devicon", "apacheairflow/apacheairflow-original", None, ["airflow"]),
    "Google Cloud":    ("devicon", "googlecloud/googlecloud-original", None, ["google cloud", "gcp", "dataflow"]),
    "Elasticsearch":   ("devicon", "elasticsearch/elasticsearch-original", None, ["elasticsearch"]),
    "Logstash":        ("devicon", "logstash/logstash-original", None, ["logstash"]),
    "Kibana":          ("devicon", "kibana/kibana-original", None, ["kibana"]),
    "Pandas":          ("devicon", "pandas/pandas-original", None, ["pandas"]),
    "Python":          ("devicon", "python/python-original", None, ["python", "pythonic"]),
    "MySQL":           ("devicon", "mysql/mysql-original", None, ["mysql"]),
    "GitHub Actions":  ("devicon", "githubactions/githubactions-original", None, ["github actions"]),
    "Git":             ("devicon", "git/git-original", None, ["git flow"]),
    "Linux":           ("devicon", "linux/linux-original", None, ["linux"]),
    "Vagrant":         ("devicon", "vagrant/vagrant-original", None, ["vagrant"]),
    # Simple Icons fallbacks (official brand color) — not in Devicon
    "Databricks":      ("simpleicons", "databricks", "FF3621", ["databricks"]),
    "MLflow":          ("simpleicons", "mlflow", "0194E2", ["mlflow"]),
    "Apache NiFi":     ("simpleicons", "apachenifi", "728E9B", ["nifi"]),
    "Splunk":          ("simpleicons", "splunk", "119DA4", ["splunk"]),
    "BigQuery":        ("simpleicons", "googlebigquery", "669DF6", ["bigquery"]),
    "Streamlit":       ("simpleicons", "streamlit", "FF4B4B", ["streamlit"]),
    "Plotly":          ("simpleicons", "plotly", "3F4F75", ["plotly"]),

    # --- JD-driven (target-role tech from lyeh-infra jd-tracker.md) ---
    # Devicon (multi-color)
    "Apache Spark":    ("devicon", "apachespark/apachespark-original", None, ["spark", "pyspark"]),
    "Apache Kafka":    ("devicon", "apachekafka/apachekafka-original", None, ["kafka"]),
    "Apache Flink":    ("simpleicons", "apacheflink", "E6526F", ["flink"]),
    "Azure":           ("devicon", "azure/azure-original", None, ["azure"]),
    "Terraform":       ("devicon", "terraform/terraform-original", None, ["terraform"]),
    "Helm":            ("devicon", "helm/helm-original", None, ["helm"]),
    "TypeScript":      ("devicon", "typescript/typescript-original", None, ["typescript"]),
    "React":           ("devicon", "react/react-original", None, ["react"]),
    "Go":              ("devicon", "go/go-original", None, ["go/", "go)", "golang"]),
    "Rust":            ("devicon", "rust/rust-original", None, ["rust"]),
    "Scala":           ("devicon", "scala/scala-original", None, ["scala"]),
    "C++":             ("devicon", "cplusplus/cplusplus-original", None, ["c++"]),
    "TensorFlow":      ("devicon", "tensorflow/tensorflow-original", None, ["tensorflow", "/tf)"]),
    "Django":          ("devicon", "django/django-plain", None, ["django"]),
    "pytest":          ("devicon", "pytest/pytest-original", None, ["pytest"]),
    "NumPy":           ("devicon", "numpy/numpy-original", None, ["numpy"]),
    "SciPy":           ("simpleicons", "scipy", "8CAAE6", ["scipy"]),
    "Neo4j":           ("devicon", "neo4j/neo4j-original", None, ["neo4j"]),
    "AWS":             ("devicon", "amazonwebservices/amazonwebservices-original-wordmark", None, ["aws", " eks", "sagemaker"]),
    # Simple Icons (official color)
    "Argo CD":         ("simpleicons", "argo", "EF7B4D", ["argocd", "argo"]),
    "LangChain":       ("simpleicons", "langchain", "1C3C3C", ["langchain", "langgraph"]),
    "OpenTelemetry":   ("simpleicons", "opentelemetry", "425CC7", ["opentelemetry", "otel"]),
    "ONNX":            ("simpleicons", "onnx", "005CED", ["onnx"]),
    "MQTT":            ("simpleicons", "mqtt", "660066", ["mqtt"]),
    "Dynatrace":       ("simpleicons", "dynatrace", "1496FF", ["dynatrace"]),
    "Palantir":        ("simpleicons", "palantir", "101113", ["palantir", "foundry"]),
    "Claude":          ("simpleicons", "claude", "D97757", ["claude", "codex"]),
}


def _slugify(name: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", name.lower()).strip("-")


def _get(url: str) -> tuple[bytes, str]:
    req = urllib.request.Request(url, headers={"User-Agent": "slidestyle-assets"})
    with urllib.request.urlopen(req, timeout=30) as r:
        return r.read(), r.headers.get("Content-Type", "")


def sync(dest: Path = ASSETS) -> dict:
    """Fetch resume + photo + icons from the registry into `dest`; write the manifest."""
    dest = Path(dest)
    (dest / "icons").mkdir(parents=True, exist_ok=True)

    resume_bytes, _ = _get(REGISTRY_URL)
    resume = json.loads(resume_bytes)
    (dest / "resume.json").write_bytes(resume_bytes)

    # --- profile photo (bump gravatar size for crisp slide use) ---
    img_url = resume.get("basics", {}).get("image", "")
    photo_rel = None
    if img_url:
        img_url = re.sub(r"size=\d+", "size=512", img_url)
        data, ctype = _get(img_url)
        ext = "jpg" if "jpeg" in ctype or "jpg" in ctype else "png"
        photo_rel = f"photo.{ext}"
        (dest / photo_rel).write_bytes(data)

    # --- tool icons: for tools that appear in the CV or the (private) JD tracker ---
    corpus = " ".join(
        kw for grp in resume.get("skills", []) for kw in grp.get("keywords", [])
    ).lower()
    jd_used = JD_DOC.exists()
    if jd_used:
        corpus += " " + JD_DOC.read_text(encoding="utf-8", errors="ignore").lower()
    keywords = corpus

    icons, skipped = {}, []
    for name, (src, slug, color, aliases) in ICON_MAP.items():
        if not any(a in keywords for a in aliases):
            continue
        url = (_DEVICON.format(slug=slug) if src == "devicon"
               else _SIMPLE.format(slug=slug, color=color))
        try:
            svg, _ = _get(url)
        except Exception as e:                       # noqa: BLE001 — best-effort per icon
            skipped.append(f"{name} ({e})")
            continue
        head = svg[:200].lstrip()
        if not (head.startswith(b"<svg") or head.startswith(b"<?xml")):
            skipped.append(f"{name} (not SVG)")
            continue
        rel = f"icons/{_slugify(name)}.svg"
        (dest / rel).write_bytes(svg)
        icons[name] = rel

    manifest = {
        "source": REGISTRY_URL,
        "jd_tools": jd_used,   # icons also seeded from the private lyeh-infra JD tracker
        "name": resume.get("basics", {}).get("name", ""),
        "photo": photo_rel,
        "icons": dict(sorted(icons.items())),
    }
    MANIFEST.write_text(json.dumps(manifest, indent=2, ensure_ascii=False) + "\n",
                        encoding="utf-8")
    if skipped:
        print("skipped:", ", ".join(skipped))
    print(f"synced {len(icons)} icons + "
          f"{'photo' if photo_rel else 'no photo'} -> {dest}")
    return manifest


# --------------------------------------------------------------------------- #
# read side — everything below works offline from the cached manifest
# --------------------------------------------------------------------------- #
_cache: dict | None = None


def manifest() -> dict:
    global _cache
    if _cache is None:
        if not MANIFEST.exists():
            raise FileNotFoundError(
                "assets/manifest.json missing — run `python -m slidestyle.assets` first")
        _cache = json.loads(MANIFEST.read_text(encoding="utf-8"))
    return _cache


def photo_path() -> Path:
    m = manifest()
    if not m.get("photo"):
        raise FileNotFoundError("no cached photo in manifest")
    return ASSETS / m["photo"]


def icons() -> dict:
    """{tool name: absolute svg Path} for every cached icon."""
    return {k: ASSETS / v for k, v in manifest()["icons"].items()}


def icon_path(name: str) -> Path:
    m = manifest()
    if name not in m["icons"]:
        raise KeyError(f"no icon for {name!r}; have: {', '.join(m['icons'])}")
    return ASSETS / m["icons"][name]


def icon_data_uri(name: str) -> str:
    b64 = base64.b64encode(icon_path(name).read_bytes()).decode("ascii")
    return f"data:image/svg+xml;base64,{b64}"


def icon_image(name: str, x: float, y: float, size: float) -> str:
    """SVG `<image>` fragment (self-contained data URI) to drop into a diagram node.
    Rasterizes correctly through cairosvg and renders in browsers."""
    uri = icon_data_uri(name)
    return (f'<image x="{x}" y="{y}" width="{size}" height="{size}" '
            f'href="{uri}" xlink:href="{uri}" '
            f'preserveAspectRatio="xMidYMid meet"/>')


# --------------------------------------------------------------------------- #
# pptx helpers (native slides)
# --------------------------------------------------------------------------- #
def add_photo(slide, left, top, width, height=None):
    """Place the profile photo on a slide (Inches). Returns the picture shape."""
    from pptx.util import Inches
    kw = {} if height is None else {"height": Inches(height)}
    return slide.shapes.add_picture(str(photo_path()), Inches(left), Inches(top),
                                    width=Inches(width), **kw)


def add_icon(slide, name, left, top, size):
    """Place a tool icon on a slide (Inches, square). Rasterizes the SVG via cairosvg
    (the `[svg]` extra) so it embeds as a crisp PNG."""
    import io
    import cairosvg
    from pptx.util import Inches
    png = cairosvg.svg2png(url=str(icon_path(name)), output_width=512, output_height=512)
    return slide.shapes.add_picture(io.BytesIO(png), Inches(left), Inches(top),
                                    width=Inches(size), height=Inches(size))


if __name__ == "__main__":
    sync()
