"""The 12 archetype builders. Each draws onto a blank slide."""
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from . import tokens as t
from . import text, shapes

def _accent_bar(slide, left, top, accent="blue", w=0.72, h=0.07):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top), Inches(w), Inches(h))
    sp.adjustments[0] = 0.5
    sp.fill.solid(); sp.fill.fore_color.rgb = t.ramp(accent, "c600")
    sp.line.fill.background()
    return sp

def _full_bg(slide, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, t.SLIDE_W, t.SLIDE_H)
    sp.fill.solid(); sp.fill.fore_color.rgb = color; sp.line.fill.background()
    slide.shapes._spTree.remove(sp._element); slide.shapes._spTree.insert(2, sp._element)
    return sp

def cover(slide, accent="blue"):
    _accent_bar(slide, 0.9, 1.5, accent)
    text.kicker(slide, 0.9, 1.75, "Product Vision · 2026", accent)
    text.title(slide, 0.86, 2.2, 11, "Designing with clarity", "cover")
    text.body(slide, 0.9, 4.1, 9, "A minimalist system for professional decks")

def section(slide, accent="blue"):
    _full_bg(slide, t.ramp(accent, "c900"))
    text.textbox(slide, (0.9, 2.6, 3, 1), "02", size=20, bold=True, color=t.ramp(accent, "c300"))
    text.textbox(slide, (0.9, 3.1, 11, 1.6), "The Approach", size=40, bold=True, color=t.WHITE, tracking=-0.4)

def statement(slide, accent="blue"):
    text.textbox(slide, (1.2, 2.6, 11, 2.2),
                 "Professionalism comes from hierarchy, not decoration.",
                 size=t.TYPE["statement"][0], bold=True, color=t.INK,
                 align=t.ALIGN_CENTER, line=1.15, tracking=-0.4)

def heading_body(slide, accent="blue"):
    text.kicker(slide, 0.9, 0.9, "Overview", accent)
    text.title(slide, 0.86, 1.35, 11, "Heading sits here", "heading")
    text.body(slide, 0.9, 2.6, 7.5,
              "Body copy in Graphite, generous line-height. One idea per slide; "
              "keep supporting detail tight and scannable.")

def bullet_list(slide, accent="blue"):
    text.kicker(slide, 0.9, 0.9, "Principles", accent)
    text.title(slide, 0.86, 1.35, 11, "What guides the design", "heading")
    items = ["Hierarchy over decoration", "One idea per slide",
             "A single accent, used deliberately", "Whitespace is a feature"]
    y = 2.7
    for it in items:
        m = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(y+0.09),
                                   Inches(0.12), Inches(0.12))
        m.fill.solid(); m.fill.fore_color.rgb = t.ramp(accent, "c600"); m.line.fill.background()
        text.textbox(slide, (1.3, y, 10, 0.6), it, size=18, color=t.GRAPHITE)
        y += 0.75

def two_column(slide, accent="blue"):
    text.kicker(slide, 0.9, 0.9, "Comparison", accent)
    text.title(slide, 0.86, 1.35, 11, "Two columns", "heading")
    text.textbox(slide, (0.9, 2.7, 5.4, 0.5), "Left", size=18, bold=True, color=t.INK)
    text.body(slide, 0.9, 3.2, 5.4, "First column body content sits here.")
    text.textbox(slide, (7.0, 2.7, 5.4, 0.5), "Right", size=18, bold=True, color=t.INK)
    text.body(slide, 7.0, 3.2, 5.4, "Second column body content sits here.")

def color_cards(slide, accent="blue"):
    text.kicker(slide, 0.9, 0.9, "Metrics", accent)
    text.title(slide, 0.86, 1.35, 11, "Impact at a glance", "heading")
    data = [("3.2×", "faster to draft", accent), ("100%", "on-brand", None), ("12", "layouts", accent)]
    x = 0.9
    for number, label, acc in data:
        shapes.kpi_tile(slide, (x, 2.9, 3.6, 1.7), number, label, accent=acc)
        x += 3.9

def highlight(slide, accent="blue"):
    text.kicker(slide, 0.9, 0.9, "Key point", accent)
    text.title(slide, 0.86, 1.35, 11, "Highlight a term", "heading")
    # tint highlight block behind a key phrase
    hl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.0),
                                Inches(5.2), Inches(1.0))
    hl.adjustments[0] = 0.12
    hl.fill.solid(); hl.fill.fore_color.rgb = t.ramp(accent, "c300"); hl.line.fill.background()
    hl.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = hl.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "the one number that matters"
    r.font.name = "Helvetica Neue"; r.font.size = Pt(20); r.font.bold = True
    r.font.color.rgb = t.ramp(accent, "c900")
    text.body(slide, 0.9, 4.3, 10, "Supporting sentence explains why it matters.")

def arch_flow(slide, accent="blue"):
    """Pattern A — horizontal flow, one focus node."""
    text.kicker(slide, 0.9, 0.9, "Architecture", accent)
    text.title(slide, 0.86, 1.35, 11, "How a deck is generated", "heading")
    nodes = [("Content", "title · body · concept", "plain"),
             ("Style template", "theme + 12 layouts", "focus"),
             ("Claude", "python-pptx", "plain"),
             (".pptx deck", "editable", "plain")]
    x, y, w, h = 0.9, 3.1, 2.5, 1.2
    prev = None
    for label, sub, kind in nodes:
        shapes.box(slide, (x, y, w, h), label, kind=kind, accent=accent, sub=sub)
        if prev is not None:
            shapes.arrow(slide, prev, y + h/2, x, y + h/2)
        prev = x + w
        x += w + 0.55

def arch_layered(slide, accent="blue"):
    """Pattern B — layered stack, three accents as categories."""
    text.kicker(slide, 0.9, 0.9, "System", accent)
    text.title(slide, 0.86, 1.35, 11, "Platform layers", "heading")
    rows = [("Presentation · UI", "Serve", "orange"),
            ("Application · services & logic", "Process", "teal"),
            ("Data · storage & ingest", "Ingest", "blue")]
    y = 2.8
    for label, tag, acc in rows:
        shapes.band(slide, (0.9, y, 11.5, 1.05), label, tag, acc)
        y += 1.25

def arch_nested(slide, accent="blue"):
    """Pattern C — nested container holding sub-nodes."""
    text.kicker(slide, 0.9, 0.9, "Architecture", accent)
    text.title(slide, 0.86, 1.35, 11, "Service boundary", "heading")
    cont = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(2.8),
                                  Inches(11.5), Inches(2.6))
    cont.adjustments[0] = 0.05
    cont.fill.solid(); cont.fill.fore_color.rgb = t.MIST; cont.line.color.rgb = t.HAIRLINE
    cont.line.width = Pt(1.5)
    text.textbox(slide, (1.15, 3.0, 6, 0.4), "CORE SERVICE", size=12, bold=True,
                 color=t.SLATE, upper=True, tracking=t.KICKER_TRACK_PT)
    subs = [("API gateway", "plain"), ("Engine", "focus"), ("Cache", "plain")]
    x = 1.2
    for label, kind in subs:
        shapes.box(slide, (x, 3.6, 3.4, 1.3), label, kind=kind, accent=accent)
        x += 3.7

def image_caption(slide, accent="blue"):
    text.kicker(slide, 0.9, 0.9, "Visual", accent)
    text.title(slide, 0.86, 1.35, 6, "Image + caption", "heading")
    ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3),
                                Inches(5.6), Inches(4.6))
    ph.adjustments[0] = 0.03
    ph.fill.solid(); ph.fill.fore_color.rgb = t.MIST; ph.line.color.rgb = t.HAIRLINE
    ph.line.width = Pt(1.5)
    p = ph.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "[ image ]"; r.font.name = "Helvetica Neue"
    r.font.size = Pt(14); r.font.color.rgb = t.SILVER
    text.body(slide, 0.9, 2.7, 5.4, "Body sits left; image fills the right half.")
    text.caption(slide, 6.8, 6.0, 5.6, "Figure 1 — caption in Slate")

def quote(slide, accent="blue"):
    text.textbox(slide, (1.2, 2.4, 11, 2),
                 "“Good design is as little design as possible.”",
                 size=36, bold=True, color=t.INK, align=t.ALIGN_CENTER, line=1.2, tracking=-0.4)
    text.textbox(slide, (1.2, 4.4, 11, 0.5), "— Dieter Rams",
                 size=16, color=t.SLATE, align=t.ALIGN_CENTER)

def closing(slide, accent="blue"):
    _accent_bar(slide, 6.13, 2.7, accent)  # roughly centered
    text.textbox(slide, (1.2, 3.0, 11, 1), "Thank you", size=30, bold=True,
                 color=t.INK, align=t.ALIGN_CENTER)
    text.textbox(slide, (1.2, 4.0, 11, 0.5), "questions → hello@example.com",
                 size=16, color=t.SLATE, align=t.ALIGN_CENTER)


def arch_swimlane(slide, accent="blue"):
    """Layout-rich architecture example: swim-lanes across environments, a shared
    data store, grouped subsystem containers, a horizontal flow with a feedback arc,
    a decision node, specialized shapes (cylinder/document/diamond), and labeled
    cross-lane connectors. Echoes a hand-made MLOps lifecycle diagram."""
    DEV, PROD = "orange", "blue"  # Dev = warm, Staging/Prod = cool (categorical)
    text.kicker(slide, 0.9, 0.42, "Architecture", accent)
    text.textbox(slide, (0.86, 0.8, 12, 0.7), "ML lifecycle across environments",
                 size=24, bold=True, color=t.INK, tracking=-0.4)

    divider_y = 4.15
    shapes._seg(slide, 0.1, divider_y, 13.23, divider_y, t.HAIRLINE, width=1.0)

    # --- containers (Mist) cover the divider where they sit ---
    shapes.container(slide, (1.95, 1.62, 9.5, 2.3), "ML Training Pipeline", accent=DEV)
    shapes.container(slide, (1.95, 4.42, 9.5, 2.35), "Staging / Prod Pipeline", accent=PROD)

    # --- shared data store straddling the divider, in the left gutter ---
    shapes.cylinder(slide, (0.52, 3.42, 1.2, 1.48), "Data\nWarehouse")

    # --- lane labels (rotated) ---
    shapes.lane_label(slide, 2.77, "Dev")
    shapes.lane_label(slide, 5.6, "Staging / Prod")

    # --- Dev lane: 3-stage flow + feedback arc + artifact ---
    dev_y, bw, bh = 2.5, 2.05, 0.8
    d1, d2, d3 = 2.35, 5.05, 7.75
    shapes.box(slide, (d1, dev_y, bw, bh), "Data Prep", kind="focus", accent=DEV)
    shapes.box(slide, (d2, dev_y, bw, bh), "Model Training", kind="focus", accent=DEV)
    shapes.box(slide, (d3, dev_y, bw, bh), "Model Validation", kind="focus", accent=DEV)
    cy = dev_y + bh / 2
    shapes._seg(slide, d1 + bw, cy, d2, cy, shapes._LINE, arrow=True)
    shapes._seg(slide, d2 + bw, cy, d3, cy, shapes._LINE, arrow=True)
    shapes.feedback_arc(slide, d3 + bw / 2, d2 + bw / 2, dev_y, "Hyperparameter tuning")
    shapes.document(slide, (10.2, dev_y, 1.2, bh), "Trained\nmodel", accent=DEV)
    shapes._seg(slide, d3 + bw, cy, 10.2, cy, shapes._LINE, arrow=True)

    # --- Prod lane: 3-stage flow + decision + retrain loop ---
    prod_y = 5.2
    p1, p2, p3 = 2.35, 5.05, 7.75
    shapes.box(slide, (p1, prod_y, bw, bh), "Data Prep", kind="plain")
    shapes.box(slide, (p2, prod_y, bw, bh), "Retrain", kind="plain")
    shapes.box(slide, (p3, prod_y, bw, bh), "Validate", kind="plain")
    pcy = prod_y + bh / 2
    shapes._seg(slide, p1 + bw, pcy, p2, pcy, shapes._LINE, arrow=True)
    shapes._seg(slide, p2 + bw, pcy, p3, pcy, shapes._LINE, arrow=True)
    shapes.diamond(slide, (10.0, pcy - 0.6, 1.2, 1.2), "Perf\ndrop?", accent=PROD)
    shapes._seg(slide, p3 + bw, pcy, 10.0, pcy, shapes._LINE, arrow=True)
    # retrain loop: down from diamond, left across, up into Data Prep (YES)
    dcx = 10.6
    shapes.elbow(slide, [(dcx, pcy + 0.6), (dcx, 6.55), (p1 + bw / 2, 6.55),
                         (p1 + bw / 2, prod_y + bh)])
    shapes._small_caps(slide, (dcx + 0.12, pcy + 0.62, 0.7, 0.3), "Yes",
                       t.ramp(PROD, "c800"), size=9)

    # --- orthogonal cross-lane connectors with white-chip labels ---
    shapes.elbow(slide, [(1.72, 3.6), (2.05, 3.6), (2.05, cy), (d1, cy)])
    shapes._chip(slide, 2.05, 3.35, "POC data")
    shapes.elbow(slide, [(1.72, 4.7), (2.05, 4.7), (2.05, pcy), (p1, pcy)])
    shapes._chip(slide, 2.05, 5.05, "Batch fetch")   # below the lane title, clear
    shapes.elbow(slide, [(10.8, dev_y + bh), (10.8, 4.42)])
    shapes._chip(slide, 10.8, 3.9, "Deploy")


def arch_swimlane_svg(slide, accent="blue"):
    """Publication-grade form of the swim-lane diagram: generated as SVG in our palette
    and embedded full-slide (renders as a crisp image; regenerate to edit). Falls back
    to a note if the [svg] extra (cairosvg) isn't installed, so the core build never breaks."""
    from .svg.generate import swimlane_svg
    try:
        from .svg.embed import embed_svg
        embed_svg(slide, swimlane_svg().encode(), 0, 0, 13.333, 7.5)
    except Exception:
        text.kicker(slide, 0.9, 0.9, "Architecture", accent)
        text.title(slide, 0.86, 1.35, 11, "Swim-lane — SVG", "heading")
        text.body(slide, 0.9, 2.6, 10.5,
                  "Publication-grade SVG version. Install the svg extra "
                  "(uv pip install -e '.[svg]') and rebuild to embed it here.")
