"""Assemble dist/slide-style.pptx: one labeled example slide per archetype."""
from pathlib import Path
from pptx import Presentation
from . import tokens as t
from .theme import apply_theme
from . import layouts, text

# (builder, caption) in gallery order
GALLERY = [
    (layouts.cover, "1 · Cover / Title"),
    (layouts.section, "2 · Section divider"),
    (layouts.statement, "3 · Big statement"),
    (layouts.heading_body, "4 · Heading + body"),
    (layouts.bullet_list, "5 · Minimal list"),
    (layouts.two_column, "6 · Two column"),
    (layouts.color_cards, "7 · Color cards / KPI"),
    (layouts.highlight, "8 · Highlight / callout"),
    (layouts.arch_flow, "9a · Architecture — flow"),
    (layouts.arch_layered, "9b · Architecture — layered"),
    (layouts.arch_nested, "9c · Architecture — nested"),
    (layouts.arch_swimlane, "9d · Architecture — swim-lane"),
    (layouts.image_caption, "10 · Image + caption"),
    (layouts.quote, "11 · Quote"),
    (layouts.closing, "12 · Closing"),
]

def build(path="dist/slide-style.pptx", primary=t.DEFAULT_ACCENT):
    prs = Presentation()
    apply_theme(prs, primary)
    for fn, cap in GALLERY:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fn(slide, accent=primary)
        # tiny archetype label bottom-left (not part of the design; a gallery marker)
        text.textbox(slide, (0.4, 7.05, 6, 0.35), cap, size=9, color=t.SILVER)
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return path

if __name__ == "__main__":
    print("built", build())
